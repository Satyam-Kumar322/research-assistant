"""
Generate a professional PowerPoint presentation on Retrieval-Augmented Generation (RAG).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ============================================================
# CONSTANTS & PALETTE
# ============================================================
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Dark theme palette
BG_PRIMARY = RGBColor(0x0F, 0x0F, 0x23)
BG_CARD = RGBColor(0x1A, 0x1A, 0x35)
BG_CARD_ALT = RGBColor(0x15, 0x15, 0x2E)
ACCENT_PURPLE = RGBColor(0x6C, 0x63, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xD4, 0xFF)
ACCENT_PINK = RGBColor(0xFF, 0x6B, 0x9D)
ACCENT_YELLOW = RGBColor(0xFF, 0xD1, 0x66)
ACCENT_GREEN = RGBColor(0x06, 0xD6, 0xA0)
TEXT_PRIMARY = RGBColor(0xF0, 0xF0, 0xF8)
TEXT_SECONDARY = RGBColor(0xA0, 0xA0, 0xC0)
TEXT_MUTED = RGBColor(0x70, 0x70, 0x90)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_SURFACE = RGBColor(0x12, 0x12, 0x2A)


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(1)):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    # Reduce corner rounding
    shape.adjustments[0] = 0.05
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=TEXT_PRIMARY, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=18, color=TEXT_PRIMARY,
                  bold=False, alignment=PP_ALIGN.LEFT, space_before=Pt(0),
                  space_after=Pt(4), font_name="Calibri"):
    """Add a new paragraph to an existing text frame."""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_bullet_item(text_frame, text, font_size=16, color=TEXT_SECONDARY,
                    bold_prefix="", bold_color=TEXT_PRIMARY, level=0):
    """Add a bullet point with optional bold prefix."""
    p = text_frame.add_paragraph()
    p.level = level
    p.space_before = Pt(4)
    p.space_after = Pt(4)

    if bold_prefix:
        run = p.add_run()
        run.text = bold_prefix
        run.font.size = Pt(font_size)
        run.font.color.rgb = bold_color
        run.font.bold = True
        run.font.name = "Calibri"

        run2 = p.add_run()
        run2.text = text
        run2.font.size = Pt(font_size)
        run2.font.color.rgb = color
        run2.font.name = "Calibri"
    else:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"

    # Set bullet character
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    buNone = pPr.find(qn('a:buNone'))
    if buNone is not None:
        pPr.remove(buNone)
    buChar = pPr.makeelement(qn('a:buChar'), {'char': '●'})
    pPr.append(buChar)
    buClr = pPr.makeelement(qn('a:buClr'), {})
    srgbClr = buClr.makeelement(qn('a:srgbClr'), {'val': '6C63FF'})
    buClr.append(srgbClr)
    pPr.append(buClr)
    buSzPct = pPr.makeelement(qn('a:buSzPct'), {'val': '70000'})
    pPr.append(buSzPct)

    return p


def add_accent_line(slide, left, top, width, color=ACCENT_PURPLE):
    """Add a thin accent line / divider."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_section_label(slide, left, top, text, color=ACCENT_CYAN):
    """Add section label like '01 · Introduction'."""
    return add_textbox(slide, left, top, Inches(6), Inches(0.4),
                       text, font_size=13, color=color, bold=True)


def add_card_with_content(slide, left, top, width, height, icon, title,
                          description, title_color=TEXT_PRIMARY,
                          border_color=None):
    """Add a glass-style card with icon, title, and description."""
    card = add_rounded_rect(slide, left, top, width, height,
                           BG_CARD, border_color or RGBColor(0x2A, 0x2A, 0x4A))

    # Icon
    add_textbox(slide, left + Inches(0.25), top + Inches(0.2),
                Inches(0.6), Inches(0.5), icon, font_size=24,
                alignment=PP_ALIGN.CENTER)

    # Title
    add_textbox(slide, left + Inches(0.25), top + Inches(0.7),
                width - Inches(0.5), Inches(0.4), title,
                font_size=15, color=title_color, bold=True)

    # Description
    add_textbox(slide, left + Inches(0.25), top + Inches(1.1),
                width - Inches(0.5), height - Inches(1.3), description,
                font_size=12, color=TEXT_MUTED)

    return card


# ============================================================
# BUILD PRESENTATION
# ============================================================
prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT
blank_layout = prs.slide_layouts[6]  # Blank layout

# ============================================================
# SLIDE 1 — TITLE
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_PRIMARY)

# Decorative shapes
circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(-1), Inches(5), Inches(5))
circle1.fill.solid()
circle1.fill.fore_color.rgb = RGBColor(0x1A, 0x18, 0x50)
circle1.line.fill.background()

circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(4), Inches(5), Inches(5))
circle2.fill.solid()
circle2.fill.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x3E)
circle2.line.fill.background()

# Top label
add_textbox(slide, Inches(0), Inches(1.0), SLIDE_WIDTH, Inches(0.4),
            "INTERN PRESENTATION  ·  JUNE 2026", font_size=13,
            color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)

# Main title
add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10.3), Inches(1.2),
            "Retrieval-Augmented Generation", font_size=46,
            color=ACCENT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(2.9), Inches(10.3), Inches(0.8),
            "(RAG)", font_size=42,
            color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)

# Subtitle
add_textbox(slide, Inches(2.5), Inches(3.8), Inches(8.3), Inches(0.7),
            "Bridging the gap between static language models and dynamic, real-world knowledge",
            font_size=18, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

# Accent line
add_accent_line(slide, Inches(5.5), Inches(4.7), Inches(2.3), ACCENT_PURPLE)

# Presenter info
add_textbox(slide, Inches(0), Inches(5.0), SLIDE_WIDTH, Inches(0.4),
            "Presented by: Satyam Kumar", font_size=18,
            color=TEXT_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0), Inches(5.4), SLIDE_WIDTH, Inches(0.4),
            "Under the guidance of Prof. (Dr.) Deepak Kumar", font_size=15,
            color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

# Date badge
badge = add_rounded_rect(slide, Inches(4.8), Inches(6.1), Inches(3.7), Inches(0.5),
                         BG_CARD, RGBColor(0x2A, 0x2A, 0x4A))
add_textbox(slide, Inches(4.8), Inches(6.15), Inches(3.7), Inches(0.4),
            "📅  19 June 2026   ·   ⏱  30 minutes", font_size=13,
            color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2 — AGENDA
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_PRIMARY)

add_section_label(slide, Inches(0.8), Inches(0.5), "OVERVIEW")
add_textbox(slide, Inches(0.8), Inches(0.9), Inches(8), Inches(0.7),
            "Agenda", font_size=38, color=TEXT_PRIMARY, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.6), Inches(1.2))

# Three agenda cards
cards_data = [
    ("📘", "Foundation", "Introduction to RAG\nWhy LLMs need RAG\nRAG Architecture & Workflow", ACCENT_PURPLE),
    ("⚙️", "Mechanics", "Document Processing & Chunking\nEmbeddings & Vector Databases\nRetrieval Mechanisms", ACCENT_CYAN),
    ("🚀", "Application", "Prompt Construction\nAdvantages & Limitations\nUse Cases & Demo", ACCENT_PINK),
]

for i, (icon, title, desc, color) in enumerate(cards_data):
    x = Inches(0.8 + i * 4.1)
    card = add_rounded_rect(slide, x, Inches(2.1), Inches(3.7), Inches(4.5),
                           BG_CARD, RGBColor(0x2A, 0x2A, 0x4A))

    # Icon circle
    icon_bg = add_rounded_rect(slide, x + Inches(0.3), Inches(2.4),
                               Inches(0.8), Inches(0.8), BG_CARD_ALT, color)
    add_textbox(slide, x + Inches(0.3), Inches(2.45), Inches(0.8), Inches(0.7),
                icon, font_size=28, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, x + Inches(0.3), Inches(3.4), Inches(3.1), Inches(0.5),
                title, font_size=22, color=color, bold=True)

    # Description lines
    txBox = slide.shapes.add_textbox(x + Inches(0.3), Inches(4.0), Inches(3.1), Inches(2.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for j, line in enumerate(desc.split('\n')):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_SECONDARY
        p.font.name = "Calibri"
        p.space_before = Pt(6)
        p.space_after = Pt(6)


# ============================================================
# SLIDE 3 — INTRODUCTION TO RAG
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_PRIMARY)

add_section_label(slide, Inches(0.8), Inches(0.5), "01  ·  INTRODUCTION")
add_textbox(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
            "What is RAG?", font_size=38, color=TEXT_PRIMARY, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.6), Inches(1.2))

# Left column - description
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(6.5), Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
run = p.add_run()
run.text = "Retrieval-Augmented Generation (RAG) "
run.font.size = Pt(16)
run.font.color.rgb = TEXT_PRIMARY
run.font.bold = True
run.font.name = "Calibri"

run2 = p.add_run()
run2.text = "is an AI framework that enhances Large Language Models by retrieving relevant information from external knowledge sources before generating responses."
run2.font.size = Pt(16)
run2.font.color.rgb = TEXT_SECONDARY
run2.font.name = "Calibri"

p2 = tf.add_paragraph()
p2.space_before = Pt(16)
run3 = p2.add_run()
run3.text = "Introduced by Lewis et al. (2020) "
run3.font.size = Pt(16)
run3.font.color.rgb = ACCENT_CYAN
run3.font.bold = True
run3.font.name = "Calibri"

run4 = p2.add_run()
run4.text = "at Facebook AI Research, RAG combines the strengths of:"
run4.font.size = Pt(16)
run4.font.color.rgb = TEXT_SECONDARY
run4.font.name = "Calibri"

add_bullet_item(tf, " — Precise access to factual data", bold_prefix="Retrieval Systems")
add_bullet_item(tf, " — Fluent, coherent language output", bold_prefix="Generative Models")

# Right column - visual card
card = add_rounded_rect(slide, Inches(8.0), Inches(2.0), Inches(4.5), Inches(4.5),
                        BG_CARD, RGBColor(0x2A, 0x2A, 0x4A))

add_textbox(slide, Inches(8.0), Inches(2.4), Inches(4.5), Inches(0.7),
            "🧠  +  📚", font_size=36, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(8.0), Inches(3.3), Inches(4.5), Inches(0.5),
            "LLM + External Knowledge", font_size=20,
            color=ACCENT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(8.3), Inches(3.9), Inches(3.9), Inches(0.5),
            "= Grounded, accurate, and up-to-date AI responses",
            font_size=14, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

# Tags
tags = [("Factual", ACCENT_CYAN), ("Current", ACCENT_GREEN), ("Verifiable", ACCENT_PURPLE)]
for i, (tag_text, tag_color) in enumerate(tags):
    tx = Inches(8.6 + i * 1.3)
    tag_shape = add_rounded_rect(slide, tx, Inches(4.7), Inches(1.1), Inches(0.4),
                                 BG_CARD_ALT, tag_color, Pt(1.5))
    add_textbox(slide, tx, Inches(4.72), Inches(1.1), Inches(0.35),
                tag_text, font_size=11, color=tag_color, bold=True,
                alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 4 — NEED FOR RAG
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_PRIMARY)

add_section_label(slide, Inches(0.8), Inches(0.5), "02  ·  MOTIVATION")
add_textbox(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
            "Why Do LLMs Need RAG?", font_size=38, color=TEXT_PRIMARY, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.6), Inches(1.2))

# Four problem cards
problems = [
    ("📅", "Knowledge Cutoff", "LLMs only know data up to their training date — no awareness of new facts or events", ACCENT_CYAN),
    ("🌀", "Hallucinations", "Models confidently generate false or fabricated information that sounds plausible", ACCENT_PINK),
    ("🏢", "Domain-Specific Data", "Enterprise & private data is not present in the public training corpus", ACCENT_YELLOW),
    ("🔍", "Source Attribution", "Need for verifiable, traceable answers with proper citations", ACCENT_GREEN),
]

for i, (icon, title, desc, color) in enumerate(problems):
    x = Inches(0.8 + i * 3.1)
    card = add_rounded_rect(slide, x, Inches(2.0), Inches(2.8), Inches(3.2),
                           BG_CARD, RGBColor(0x2A, 0x2A, 0x4A))

    add_textbox(slide, x, Inches(2.2), Inches(2.8), Inches(0.6),
                icon, font_size=32, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, x + Inches(0.2), Inches(2.9), Inches(2.4), Inches(0.4),
                title, font_size=15, color=color, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, x + Inches(0.2), Inches(3.4), Inches(2.4), Inches(1.5),
                desc, font_size=12, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# Bottom warning card
warn_card = add_rounded_rect(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.3),
                             RGBColor(0x1A, 0x12, 0x20), ACCENT_PINK, Pt(1))

txBox = slide.shapes.add_textbox(Inches(1.2), Inches(5.7), Inches(11), Inches(0.9))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "⚠️  Without RAG, "
run.font.size = Pt(15)
run.font.color.rgb = ACCENT_PINK
run.font.bold = True
run.font.name = "Calibri"
run2 = p.add_run()
run2.text = "fine-tuning on new data is expensive, slow, and risks catastrophic forgetting."
run2.font.size = Pt(15)
run2.font.color.rgb = ACCENT_PINK
run2.font.name = "Calibri"
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
run3 = p2.add_run()
run3.text = "RAG offers a lightweight, dynamic alternative — no retraining needed."
run3.font.size = Pt(13)
run3.font.color.rgb = TEXT_MUTED
run3.font.name = "Calibri"


# ============================================================
# SLIDE 5 — RAG ARCHITECTURE & WORKFLOW
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_PRIMARY)

add_section_label(slide, Inches(0.8), Inches(0.5), "03  ·  ARCHITECTURE")
add_textbox(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
            "RAG Architecture & Workflow", font_size=38, color=TEXT_PRIMARY, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.6), Inches(1.2))

# Flow diagram
flow_steps = [
    ("👤", "User\nQuery", ACCENT_CYAN),
    ("🔢", "Embed\nQuery", ACCENT_PURPLE),
    ("🔍", "Vector\nSearch", ACCENT_CYAN),
    ("📄", "Retrieve\nDocs", ACCENT_GREEN),
    ("📝", "Augment\nPrompt", ACCENT_YELLOW),
    ("🤖", "LLM\nResponse", ACCENT_PINK),
]

for i, (icon, label, color) in enumerate(flow_steps):
    x = Inches(0.6 + i * 2.15)
    step_card = add_rounded_rect(slide, x, Inches(2.0), Inches(1.5), Inches(1.4),
                                 BG_CARD, color, Pt(1.5))
    add_textbox(slide, x, Inches(2.05), Inches(1.5), Inches(0.5),
                icon, font_size=22, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(2.55), Inches(1.5), Inches(0.7),
                label, font_size=11, color=TEXT_PRIMARY, bold=True,
                alignment=PP_ALIGN.CENTER)

    if i < len(flow_steps) - 1:
        add_textbox(slide, x + Inches(1.5), Inches(2.35), Inches(0.65), Inches(0.5),
                    "→", font_size=22, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

# Two pipeline cards
# Indexing Pipeline
card1 = add_rounded_rect(slide, Inches(0.8), Inches(3.8), Inches(5.6), Inches(3.2),
                         BG_CARD, RGBColor(0x2A, 0x2A, 0x4A))
add_textbox(slide, Inches(1.1), Inches(3.95), Inches(5), Inches(0.4),
            "📥  Indexing Pipeline (Offline)", font_size=17,
            color=ACCENT_CYAN, bold=True)

txBox = slide.shapes.add_textbox(Inches(1.1), Inches(4.5), Inches(5), Inches(2.3))
tf = txBox.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = ""
add_bullet_item(tf, "Load documents (PDF, HTML, DB, API)")
add_bullet_item(tf, "Chunk into manageable segments")
add_bullet_item(tf, "Generate embeddings via encoder model")
add_bullet_item(tf, "Store in vector database")

# Query Pipeline
card2 = add_rounded_rect(slide, Inches(6.9), Inches(3.8), Inches(5.6), Inches(3.2),
                         BG_CARD, RGBColor(0x2A, 0x2A, 0x4A))
add_textbox(slide, Inches(7.2), Inches(3.95), Inches(5), Inches(0.4),
            "📤  Query Pipeline (Online)", font_size=17,
            color=ACCENT_PINK, bold=True)

txBox = slide.shapes.add_textbox(Inches(7.2), Inches(4.5), Inches(5), Inches(2.3))
tf = txBox.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = ""
add_bullet_item(tf, "Embed user query with same model")
add_bullet_item(tf, "Retrieve top-k relevant chunks")
add_bullet_item(tf, "Construct augmented prompt with context")
add_bullet_item(tf, "Generate grounded response via LLM")


# ============================================================
# SLIDE 6 — DOCUMENT PROCESSING & CHUNKING
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_PRIMARY)

add_section_label(slide, Inches(0.8), Inches(0.5), "04  ·  DATA INGESTION")
add_textbox(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
            "Document Processing & Chunking", font_size=38, color=TEXT_PRIMARY, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.6), Inches(1.2))

# Left: Document Processing
add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(0.4),
            "Document Processing", font_size=20, color=TEXT_PRIMARY, bold=True)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(5.5), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = ""
add_bullet_item(tf, " — Parse PDFs, DOCX, HTML, Markdown, CSVs", bold_prefix="Ingestion")
add_bullet_item(tf, " — Remove noise, headers, footers, formatting", bold_prefix="Cleaning")
add_bullet_item(tf, " — Title, date, author, section tags", bold_prefix="Metadata Extraction")
add_bullet_item(tf, " — Unicode, case folding, deduplication", bold_prefix="Normalization")

# Right: Chunking Strategies
add_textbox(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(0.4),
            "Chunking Strategies", font_size=20, color=TEXT_PRIMARY, bold=True)

chunks_data = [
    ("Fixed-Size Chunking", "Split every N tokens/characters. Simple but may break context.", ACCENT_PURPLE),
    ("Recursive / Semantic Chunking", "Split by paragraph, sentence, then token — preserves meaning.", ACCENT_CYAN),
    ("Overlapping Windows", "Chunks overlap by 10-20% to retain boundary context.", ACCENT_GREEN),
]

for i, (title, desc, color) in enumerate(chunks_data):
    y = Inches(2.5 + i * 1.15)
    # Left border card
    border_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(7.0), y, Pt(5), Inches(0.95)
    )
    border_line.fill.solid()
    border_line.fill.fore_color.rgb = color
    border_line.line.fill.background()

    card = add_rounded_rect(slide, Inches(7.1), y, Inches(5.2), Inches(0.95),
                           BG_CARD, RGBColor(0x2A, 0x2A, 0x4A))

    add_textbox(slide, Inches(7.3), y + Inches(0.1), Inches(4.8), Inches(0.35),
                title, font_size=14, color=color, bold=True)
    add_textbox(slide, Inches(7.3), y + Inches(0.5), Inches(4.8), Inches(0.4),
                desc, font_size=12, color=TEXT_MUTED)

# Bottom tip
tip_card = add_rounded_rect(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.9),
                            RGBColor(0x0D, 0x1E, 0x1E), ACCENT_GREEN, Pt(1))
add_textbox(slide, Inches(1.2), Inches(5.95), Inches(11), Inches(0.6),
            "💡 Best Practice: Typical chunk size = 256–1024 tokens with 10–20% overlap for optimal retrieval recall.",
            font_size=14, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 7 — EMBEDDINGS & VECTOR DATABASES
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_PRIMARY)

add_section_label(slide, Inches(0.8), Inches(0.5), "05  ·  EMBEDDING & STORAGE")
add_textbox(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
            "Embeddings & Vector Databases", font_size=38, color=TEXT_PRIMARY, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.6), Inches(1.2))

# Left: Embeddings
card1 = add_rounded_rect(slide, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.8),
                         BG_CARD, RGBColor(0x2A, 0x2A, 0x4A))

add_textbox(slide, Inches(1.1), Inches(2.2), Inches(5), Inches(0.4),
            "🔢  Embeddings", font_size=19, color=ACCENT_PURPLE, bold=True)

txBox = slide.shapes.add_textbox(Inches(1.1), Inches(2.7), Inches(5), Inches(1.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Dense vector representations that capture semantic meaning of text in high-dimensional space."
p.font.size = Pt(14)
p.font.color.rgb = TEXT_SECONDARY
p.font.name = "Calibri"

# Embedding model tags
models = [
    ("OpenAI text-embedding-3", ACCENT_CYAN), ("Sentence-BERT", ACCENT_GREEN),
    ("Cohere Embed v3", ACCENT_PURPLE), ("BGE / E5", ACCENT_YELLOW)
]
for i, (model, color) in enumerate(models):
    x = Inches(1.1 + (i % 2) * 2.5)
    y = Inches(3.8 + (i // 2) * 0.5)
    tag_shape = add_rounded_rect(slide, x, y, Inches(2.2), Inches(0.38),
                                 BG_CARD_ALT, color, Pt(1))
    add_textbox(slide, x, y + Inches(0.02), Inches(2.2), Inches(0.33),
                model, font_size=11, color=color, bold=True,
                alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.1), Inches(5.0), Inches(5), Inches(0.6),
            "Similarity measured via cosine similarity or dot product between embedding vectors.",
            font_size=13, color=TEXT_MUTED)

# Right: Vector Databases
card2 = add_rounded_rect(slide, Inches(6.9), Inches(2.0), Inches(5.6), Inches(4.8),
                         BG_CARD, RGBColor(0x2A, 0x2A, 0x4A))

add_textbox(slide, Inches(7.2), Inches(2.2), Inches(5), Inches(0.4),
            "🗄️  Vector Databases", font_size=19, color=ACCENT_GREEN, bold=True)

txBox = slide.shapes.add_textbox(Inches(7.2), Inches(2.7), Inches(5), Inches(0.7))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Specialized databases optimized for Approximate Nearest Neighbor (ANN) search at scale."
p.font.size = Pt(14)
p.font.color.rgb = TEXT_SECONDARY
p.font.name = "Calibri"

vdbs = [
    ("Pinecone", "Managed, scalable SaaS"),
    ("ChromaDB", "Lightweight, open-source"),
    ("Weaviate", "GraphQL-native, hybrid search"),
    ("FAISS", "Meta's in-memory ANN library"),
    ("Qdrant", "Rust-based, fast filtering"),
]
for i, (name, desc) in enumerate(vdbs):
    y = Inches(3.5 + i * 0.55)
    # Row
    add_textbox(slide, Inches(7.4), y, Inches(1.8), Inches(0.4),
                name, font_size=14, color=TEXT_PRIMARY, bold=True)
    add_textbox(slide, Inches(9.3), y, Inches(3.0), Inches(0.4),
                desc, font_size=12, color=TEXT_MUTED)
    if i < len(vdbs) - 1:
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(7.4), y + Inches(0.45),
            Inches(4.5), Pt(1)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x40)
        line.line.fill.background()


# ============================================================
# SLIDE 8 — RETRIEVAL MECHANISMS
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_PRIMARY)

add_section_label(slide, Inches(0.8), Inches(0.5), "06  ·  RETRIEVAL")
add_textbox(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
            "Retrieval Mechanisms", font_size=38, color=TEXT_PRIMARY, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.6), Inches(1.2))

# Three retrieval type cards
retrieval_types = [
    ("📐", "Dense Retrieval", "Semantic similarity via embeddings (DPR, Contriever). Best for meaning-based search.", "Semantic Match", ACCENT_PURPLE),
    ("🔤", "Sparse Retrieval", "Keyword-based methods like BM25, TF-IDF. Fast and interpretable for exact matching.", "Keyword Match", ACCENT_CYAN),
    ("🔀", "Hybrid Retrieval", "Combine dense + sparse with reciprocal rank fusion for best-of-both-worlds results.", "Best Accuracy", ACCENT_GREEN),
]

for i, (icon, title, desc, tag, color) in enumerate(retrieval_types):
    x = Inches(0.8 + i * 4.1)
    card = add_rounded_rect(slide, x, Inches(2.0), Inches(3.7), Inches(2.8),
                           BG_CARD, RGBColor(0x2A, 0x2A, 0x4A))

    add_textbox(slide, x, Inches(2.15), Inches(3.7), Inches(0.5),
                icon, font_size=28, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), Inches(2.7), Inches(3.3), Inches(0.4),
                title, font_size=16, color=color, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), Inches(3.2), Inches(3.3), Inches(0.9),
                desc, font_size=12, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

    # Tag
    tag_shape = add_rounded_rect(slide, x + Inches(1.1), Inches(4.2), Inches(1.5), Inches(0.35),
                                 BG_CARD_ALT, color, Pt(1))
    add_textbox(slide, x + Inches(1.1), Inches(4.22), Inches(1.5), Inches(0.3),
                tag, font_size=10, color=color, bold=True,
                alignment=PP_ALIGN.CENTER)

# Advanced Retrieval card
adv_card = add_rounded_rect(slide, Inches(0.8), Inches(5.1), Inches(11.7), Inches(2.0),
                            BG_CARD, RGBColor(0x2A, 0x2A, 0x4A))

add_textbox(slide, Inches(1.1), Inches(5.2), Inches(5), Inches(0.4),
            "🔧  Advanced Retrieval Techniques", font_size=16,
            color=ACCENT_YELLOW, bold=True)

# Left col
txBox = slide.shapes.add_textbox(Inches(1.1), Inches(5.7), Inches(5.3), Inches(1.2))
tf = txBox.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = ""
add_bullet_item(tf, " — Cross-encoder rescoring of top-k results", bold_prefix="Re-ranking", font_size=13)
add_bullet_item(tf, " — HyDE, multi-query rewriting", bold_prefix="Query Expansion", font_size=13)

# Right col
txBox = slide.shapes.add_textbox(Inches(6.8), Inches(5.7), Inches(5.3), Inches(1.2))
tf = txBox.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = ""
add_bullet_item(tf, " — Date, source, category constraints", bold_prefix="Metadata Filtering", font_size=13)
add_bullet_item(tf, " — Maximal Marginal Relevance for diversity", bold_prefix="MMR", font_size=13)


# ============================================================
# SLIDE 9 — PROMPT CONSTRUCTION
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_PRIMARY)

add_section_label(slide, Inches(0.8), Inches(0.5), "07  ·  PROMPT ENGINEERING")
add_textbox(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
            "Prompt Construction & Context Window", font_size=38, color=TEXT_PRIMARY, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.6), Inches(1.2))

# Code block
code_card = add_rounded_rect(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(3.2),
                              RGBColor(0x0A, 0x0A, 0x18), RGBColor(0x2A, 0x2A, 0x4A))

code_text = '''# RAG Prompt Template

prompt = f"""
You are a knowledgeable research assistant.
Answer the question using ONLY the provided context.
If the context doesn't contain the answer, say "I don't know."

### Context:
{retrieved_documents}

### Question:
{user_query}

### Answer:
"""'''

txBox = slide.shapes.add_textbox(Inches(1.8), Inches(2.2), Inches(9.7), Inches(2.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = code_text
p.font.size = Pt(13)
p.font.color.rgb = ACCENT_GREEN
p.font.name = "Consolas"

# Three stat cards
stats = [
    ("128K", "GPT-4o Context Window", ACCENT_PURPLE),
    ("1M", "Gemini 2.5 Pro Context", ACCENT_CYAN),
    ("Top-k", "Typically k = 3–10 chunks", ACCENT_GREEN),
]

for i, (number, label, color) in enumerate(stats):
    x = Inches(1.0 + i * 4.0)
    card = add_rounded_rect(slide, x, Inches(5.6), Inches(3.5), Inches(1.4),
                           BG_CARD, RGBColor(0x2A, 0x2A, 0x4A))
    add_textbox(slide, x, Inches(5.7), Inches(3.5), Inches(0.7),
                number, font_size=34, color=color, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(6.3), Inches(3.5), Inches(0.4),
                label, font_size=12, color=TEXT_MUTED,
                alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 10 — ADVANTAGES
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_PRIMARY)

add_section_label(slide, Inches(0.8), Inches(0.5), "08  ·  BENEFITS")
add_textbox(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
            "Advantages of RAG", font_size=38, color=TEXT_PRIMARY, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.6), Inches(1.2))

advantages = [
    ("✅  Reduced Hallucinations", "Responses grounded in retrieved factual documents, not parametric memory alone.", ACCENT_GREEN),
    ("✅  Up-to-Date Knowledge", "Update the knowledge base anytime — no expensive retraining cycle needed.", ACCENT_CYAN),
    ("✅  Cost-Effective", "Far cheaper than fine-tuning. Only requires embedding + storage infrastructure.", ACCENT_PURPLE),
    ("✅  Source Traceability", "Every answer can cite its source documents — critical for compliance and trust.", ACCENT_YELLOW),
    ("✅  Domain Adaptability", "Inject any domain knowledge (legal, medical, technical) without changing the model.", ACCENT_PINK),
    ("✅  Data Privacy", "Sensitive data stays in your vector DB — never sent for model training.", ACCENT_CYAN),
]

for i, (title, desc, color) in enumerate(advantages):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(2.0 + row * 1.55)

    # Left border
    border = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, y, Pt(5), Inches(1.3)
    )
    border.fill.solid()
    border.fill.fore_color.rgb = color
    border.line.fill.background()

    card = add_rounded_rect(slide, x + Pt(5), y, Inches(5.7), Inches(1.3),
                           BG_CARD, RGBColor(0x2A, 0x2A, 0x4A))

    add_textbox(slide, x + Inches(0.2), y + Inches(0.15), Inches(5.3), Inches(0.4),
                title, font_size=15, color=color, bold=True)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.6), Inches(5.3), Inches(0.6),
                desc, font_size=13, color=TEXT_MUTED)


# ============================================================
# SLIDE 11 — LIMITATIONS
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_PRIMARY)

add_section_label(slide, Inches(0.8), Inches(0.5), "09  ·  CHALLENGES")
add_textbox(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
            "Limitations & Challenges", font_size=38, color=TEXT_PRIMARY, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.6), Inches(1.2))

# Left: Limitations timeline
limitations = [
    ("Retrieval Quality Dependency", "Output is only as good as retrieved context — garbage in, garbage out."),
    ("Latency Overhead", "Additional embedding + search step adds 200–800ms per query."),
    ("Context Window Limits", "Too many chunks can overwhelm the model or exceed token limits."),
    ("Complex Multi-Hop Reasoning", "Single retrieval may miss answers requiring chained reasoning across documents."),
]

for i, (title, desc) in enumerate(limitations):
    y = Inches(2.1 + i * 1.15)

    # Timeline dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.15), y + Inches(0.08),
                                  Inches(0.2), Inches(0.2))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT_PURPLE
    dot.line.fill.background()

    # Timeline line
    if i < len(limitations) - 1:
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1.23), y + Inches(0.3),
            Pt(2), Inches(0.95)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_PURPLE
        line.line.fill.background()

    add_textbox(slide, Inches(1.6), y, Inches(4.5), Inches(0.35),
                title, font_size=15, color=TEXT_PRIMARY, bold=True)
    add_textbox(slide, Inches(1.6), y + Inches(0.4), Inches(4.5), Inches(0.5),
                desc, font_size=12, color=TEXT_MUTED)

# Right: RAG vs Fine-Tuning
comp_card = add_rounded_rect(slide, Inches(6.8), Inches(2.1), Inches(5.7), Inches(3.8),
                             RGBColor(0x1A, 0x18, 0x28), RGBColor(0x35, 0x30, 0x50))

add_textbox(slide, Inches(7.1), Inches(2.25), Inches(5), Inches(0.4),
            "⚡  RAG vs Fine-Tuning", font_size=17, color=ACCENT_YELLOW, bold=True)

# Table header
headers = [("Aspect", Inches(7.1)), ("RAG", Inches(9.0)), ("Fine-Tuning", Inches(10.8))]
for text, x in headers:
    add_textbox(slide, x, Inches(2.8), Inches(1.7), Inches(0.35),
                text, font_size=12, color=ACCENT_CYAN, bold=True)

# Table separator
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.1), Inches(3.15),
                               Inches(5.1), Pt(1))
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(0x30, 0x30, 0x50)
line.line.fill.background()

# Table rows
rows = [
    ("Setup Cost", "Low", "High"),
    ("Update Speed", "Real-time", "Hours/Days"),
    ("Traceability", "✅ Yes", "❌ No"),
    ("Best For", "Knowledge", "Style/Behavior"),
]
for i, (aspect, rag, ft) in enumerate(rows):
    y = Inches(3.3 + i * 0.45)
    add_textbox(slide, Inches(7.1), y, Inches(1.7), Inches(0.35),
                aspect, font_size=12, color=TEXT_MUTED)
    add_textbox(slide, Inches(9.0), y, Inches(1.7), Inches(0.35),
                rag, font_size=12, color=TEXT_SECONDARY)
    add_textbox(slide, Inches(10.8), y, Inches(1.7), Inches(0.35),
                ft, font_size=12, color=TEXT_SECONDARY)

# Bottom tip
tip_card = add_rounded_rect(slide, Inches(6.8), Inches(6.1), Inches(5.7), Inches(0.8),
                            BG_CARD, RGBColor(0x2A, 0x2A, 0x4A))
add_textbox(slide, Inches(7.1), Inches(6.2), Inches(5.1), Inches(0.5),
            "💡 Pro Tip: Use RAG + Fine-Tuning together for maximum performance.",
            font_size=13, color=ACCENT_YELLOW, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 12 — USE CASES
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_PRIMARY)

add_section_label(slide, Inches(0.8), Inches(0.5), "10  ·  APPLICATIONS")
add_textbox(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
            "Real-World Use Cases", font_size=38, color=TEXT_PRIMARY, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.6), Inches(1.2))

use_cases = [
    ("🏥", "Healthcare", "Medical literature Q&A, drug interaction lookup, clinical decision support.", ACCENT_PURPLE),
    ("⚖️", "Legal Tech", "Contract analysis, case law research, regulatory compliance checking.", ACCENT_CYAN),
    ("🎓", "Education", "AI tutors grounded in textbooks, personalized learning from course materials.", ACCENT_GREEN),
    ("💬", "Customer Support", "Enterprise chatbots trained on knowledge base articles and FAQs.", ACCENT_PINK),
    ("💻", "Developer Tools", "Code-aware assistants that search documentation and codebase in real-time.", ACCENT_YELLOW),
    ("🔬", "Research", "Scientific paper analysis, literature review, hypothesis generation.", ACCENT_PURPLE),
]

for i, (icon, title, desc, color) in enumerate(use_cases):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(2.0 + row * 2.0)

    # Left border
    border = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, y, Pt(5), Inches(1.6)
    )
    border.fill.solid()
    border.fill.fore_color.rgb = color
    border.line.fill.background()

    card = add_rounded_rect(slide, x + Pt(5), y, Inches(3.6), Inches(1.6),
                           BG_CARD, RGBColor(0x2A, 0x2A, 0x4A))

    add_textbox(slide, x + Inches(0.2), y + Inches(0.15), Inches(3.2), Inches(0.4),
                f"{icon}  {title}", font_size=16, color=color, bold=True)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.65), Inches(3.2), Inches(0.8),
                desc, font_size=12, color=TEXT_MUTED)

# Bottom highlight card
hl_card = add_rounded_rect(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.9),
                           RGBColor(0x14, 0x12, 0x30), ACCENT_PURPLE, Pt(1))
add_textbox(slide, Inches(1.2), Inches(6.35), Inches(11), Inches(0.5),
            "🚀  Our Research Assistant Project — is a working RAG system that processes documents, chunks text, and retrieves context-aware answers!",
            font_size=14, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 13 — DEMO / CASE STUDY
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_PRIMARY)

add_section_label(slide, Inches(0.8), Inches(0.5), "11  ·  DEMO")
add_textbox(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
            "Case Study: Research Assistant RAG Pipeline", font_size=34,
            color=TEXT_PRIMARY, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.6), Inches(1.2))

# Left: System Architecture code block
add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(0.4),
            "System Architecture", font_size=18, color=ACCENT_CYAN, bold=True)

code_card = add_rounded_rect(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(4.5),
                              RGBColor(0x0A, 0x0A, 0x18), RGBColor(0x2A, 0x2A, 0x4A))

arch_text = """# Our RAG Pipeline Stack

Backend:     FastAPI + Python
LLM:           Google Gemini API
Embedding:  Sentence Transformers
Database:    SQLite + Vector Store
Auth:           Firebase Authentication
Frontend:    HTML/CSS/JS Templates

# Document Processing Pipeline
1. Upload PDF/DOCX/TXT
2. Extract & clean text
3. Chunk into segments
4. Generate embeddings
5. Store for retrieval
6. Query → Retrieve → Generate"""

txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.6), Inches(5.1), Inches(4.2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = arch_text
p.font.size = Pt(12)
p.font.color.rgb = ACCENT_GREEN
p.font.name = "Consolas"

# Right: Key Features
add_textbox(slide, Inches(6.8), Inches(2.0), Inches(5.5), Inches(0.4),
            "Key Features Demonstrated", font_size=18, color=ACCENT_PINK, bold=True)

features = [
    ("📄  Multi-Format Document Upload", "Supports PDF, DOCX, TXT with automated text extraction and cleaning.", ACCENT_GREEN),
    ("🧹  Intelligent Text Cleaning", "Custom pipeline removes noise, normalizes text, preserves semantic structure.", ACCENT_PURPLE),
    ("🔍  Context-Aware Retrieval", "Semantic search finds the most relevant document chunks for each query.", ACCENT_CYAN),
    ("🤖  Grounded LLM Responses", "Gemini generates answers strictly from retrieved context with citations.", ACCENT_YELLOW),
]

for i, (title, desc, color) in enumerate(features):
    y = Inches(2.5 + i * 1.2)

    border = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.8), y, Pt(5), Inches(1.0)
    )
    border.fill.solid()
    border.fill.fore_color.rgb = color
    border.line.fill.background()

    card = add_rounded_rect(slide, Inches(6.9), y, Inches(5.4), Inches(1.0),
                           BG_CARD, RGBColor(0x2A, 0x2A, 0x4A))

    add_textbox(slide, Inches(7.1), y + Inches(0.1), Inches(5), Inches(0.35),
                title, font_size=14, color=color, bold=True)
    add_textbox(slide, Inches(7.1), y + Inches(0.5), Inches(5), Inches(0.4),
                desc, font_size=12, color=TEXT_MUTED)


# ============================================================
# SLIDE 14 — THANK YOU
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_PRIMARY)

# Decorative circles
c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(6), Inches(6))
c1.fill.solid()
c1.fill.fore_color.rgb = RGBColor(0x14, 0x12, 0x38)
c1.line.fill.background()

c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(4), Inches(5), Inches(5))
c2.fill.solid()
c2.fill.fore_color.rgb = RGBColor(0x10, 0x18, 0x30)
c2.line.fill.background()

# Big ring (outer)
ring_outer = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.4), Inches(0.8),
                                     Inches(2.5), Inches(2.5))
ring_outer.fill.solid()
ring_outer.fill.fore_color.rgb = ACCENT_PURPLE
ring_outer.line.fill.background()

# Ring inner
ring_inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.6), Inches(1.0),
                                     Inches(2.1), Inches(2.1))
ring_inner.fill.solid()
ring_inner.fill.fore_color.rgb = BG_PRIMARY
ring_inner.line.fill.background()

# Emoji in ring
add_textbox(slide, Inches(5.4), Inches(1.3), Inches(2.5), Inches(1.5),
            "🎓", font_size=48, alignment=PP_ALIGN.CENTER)

# Thank You text
add_textbox(slide, Inches(0), Inches(3.6), SLIDE_WIDTH, Inches(0.9),
            "Thank You!", font_size=44, color=TEXT_PRIMARY, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(2), Inches(4.4), Inches(9.3), Inches(0.5),
            "Questions, feedback, and discussion welcome",
            font_size=18, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5.8), Inches(5.1), Inches(1.7), ACCENT_PURPLE)

add_textbox(slide, Inches(0), Inches(5.4), SLIDE_WIDTH, Inches(0.4),
            "Satyam Kumar", font_size=20, color=TEXT_PRIMARY,
            bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0), Inches(5.8), SLIDE_WIDTH, Inches(0.4),
            "Under the guidance of Prof. (Dr.) Deepak Kumar",
            font_size=15, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# Tags at bottom
tags_bottom = [
    ("RAG", ACCENT_PURPLE), ("LLM", ACCENT_CYAN),
    ("AI/ML", ACCENT_GREEN), ("NLP", ACCENT_PINK)
]
for i, (tag_text, color) in enumerate(tags_bottom):
    x = Inches(4.8 + i * 1.1)
    tag_shape = add_rounded_rect(slide, x, Inches(6.5), Inches(0.9), Inches(0.4),
                                 BG_CARD_ALT, color, Pt(1.5))
    add_textbox(slide, x, Inches(6.52), Inches(0.9), Inches(0.35),
                tag_text, font_size=12, color=color, bold=True,
                alignment=PP_ALIGN.CENTER)


# ============================================================
# SAVE
# ============================================================
output_path = r"c:\Users\SATYAM\OneDrive\Desktop\research_assistant\RAG_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
